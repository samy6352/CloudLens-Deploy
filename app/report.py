"""Exporting the dashboard as a report someone can send to their finance team.

Four formats, chosen because they are what people actually do with cost data:

- **XLSX** — the one finance asks for. Multi-sheet, real numbers with currency formatting, so
  it can be pivoted rather than retyped.
- **PPTX** — the one that goes to a steering committee. Native PowerPoint charts, not pasted
  images, so the deck can be restyled without coming back here.
- **CSV** — the lowest common denominator; imports anywhere.
- **Markdown** — for a ticket, a wiki or a pull request description.

PDF is deliberately absent: the browser already prints this page well (there are print styles),
and a headless renderer on a Free-tier App Service would be a large dependency producing a
worse result than Ctrl-P.

Everything is generated in memory and streamed. Writing a temp file would be a needless
failure mode on a 1 GB instance, and reports are small.
"""
from __future__ import annotations

import csv
import io
import logging
from datetime import datetime, timezone
from typing import Any

from .currency import symbol as fx_symbol

log = logging.getLogger("cloudlens.report")

FORMATS = {
    "xlsx": ("Excel workbook",
             "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    "pptx": ("PowerPoint deck",
             "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    "csv": ("CSV", "text/csv"),
    "md": ("Markdown", "text/markdown"),
    "json": ("JSON", "application/json"),
}

# Which parts of an area section a report can carry. Selected individually because a deck for
# a steering committee wants the trend and nothing else, while a handover to an engineer wants
# the resource list and could do without the chart.
BLOCKS = ("trend", "services", "regions", "resources")

# Live datasets. Separate from the warehouse sections because they call Azure — including one
# turns a sub-second export into a slow one, so it has to be a deliberate choice.
LIVE = {
    "waste": "Orphaned resources",
    "rightsizing": "Rightsizing",
    "esu": "End of support",
    "advisor": "Advisor",
}


def default_selection() -> dict[str, Any]:
    """Everything from the warehouse, nothing live. The safe, fast default."""
    return {"summary": True, "sections": None, "blocks": list(BLOCKS), "live": []}


def collect(scope: list[str] | None, days: int = 30,
            selection: dict[str, Any] | None = None,
            live_data: dict[str, Any] | None = None,
            currency: str | None = None) -> dict[str, Any]:
    """Everything a report needs, filtered to what was asked for.

    Deliberately the same `Dashboard` the screen uses: a report that disagreed with the page
    it was exported from would be worse than no report. `currency` is part of that agreement —
    a report ordered while the screen showed rupees has to contain rupees.

    `selection` chooses content:
      summary  — include the KPIs and the by-area table
      sections — list of area ids, or None for all
      blocks   — which parts of each area (trend / services / regions / resources)
      live     — which live datasets to embed, fetched by the caller and passed in `live_data`
    """
    from .dashboard import get_dashboard

    sel = {**default_selection(), **(selection or {})}
    blocks = set(sel.get("blocks") or BLOCKS)
    wanted = sel.get("sections")

    d = get_dashboard()
    summary = d.sections(scope, days=days, currency=currency)
    areas = [s for s in summary.get("sections", []) if s.get("cost")]
    if wanted is not None:
        keep = set(wanted)
        areas = [a for a in areas if a["id"] in keep]

    sections = []
    for area in areas:
        try:
            section = d.section(area["id"], scope, days=days, currency=currency)
        except KeyError:  # a group with no rows in this window
            continue
        # Drop the blocks that were not asked for, rather than writing empty sheets and slides.
        if "trend" not in blocks:
            section["trend"] = {"labels": [], "values": []}
        for name in ("services", "regions"):
            if name not in blocks:
                section[name] = []
        if "resources" not in blocks:
            section["resources_top"] = []
        sections.append(section)

    return {
        "generated": datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC"),
        "days": days,
        "as_of": summary.get("as_of"),
        "currency": summary.get("currency") or "USD",
        # The headline total always describes the whole estate in scope, even when only some
        # areas are exported — otherwise a partial report looks like a smaller bill.
        "total": summary.get("total") or 0.0,
        "areas_total": round(sum(a["cost"] for a in areas), 2),
        "mixed_currency": summary.get("mixed"),
        "areas": areas if sel.get("summary") else [],
        "sections": sections,
        "live": live_data or {},
        "partial": wanted is not None or set(blocks) != set(BLOCKS) or not sel.get("summary"),
        "scoped": bool(scope),
        "subscription_count": len(scope) if scope else None,
    }


def _money_format(currency: str) -> str:
    mark = fx_symbol(currency)
    return f'"{mark}"#,##0.00' if mark else "#,##0.00"


def live_tables(data: dict[str, Any]) -> list[dict[str, Any]]:
    """Flatten whatever live datasets were included into plain titled tables.

    Each live tool returns its own shape; every export format wants the same thing — a title,
    some headline facts, column names and rows. Normalising once here means the four writers
    below stay simple and cannot disagree about what a finding looks like.
    """
    out: list[dict[str, Any]] = []
    live = data.get("live") or {}
    cur = data["currency"]

    waste = live.get("waste")
    if waste:
        for finding in waste.get("findings", []):
            if not finding.get("count"):
                continue
            out.append({
                "key": "waste",
                "title": str(finding.get("category", "finding")).replace("_", " ").title(),
                "facts": [
                    f"{finding['count']} item(s)",
                    f"{finding.get('cost') or 0:,.2f} {cur} over {waste.get('period_days', '?')} days",
                ],
                "columns": ["Resource", "Resource group", "Region", f"Cost ({cur})"],
                "rows": [[i.get("name") or i.get("id") or "—", i.get("resource_group") or "—",
                          i.get("location") or i.get("region") or "—", i.get("cost")]
                         for i in (finding.get("items") or [])],
                "money_col": 4,
            })

    rs = live.get("rightsizing")
    if rs:
        out.append({
            "key": "rightsizing",
            "title": "Rightsizing",
            "facts": [f"{rs.get('count', 0)} VM(s)", f"{rs.get('idle_count', 0)} idle",
                      f"idle cost {rs.get('idle_cost') or 0:,.2f} {cur}"],
            "columns": ["Name", "Size", "Power state", "CPU avg %", "CPU peak %", f"Cost ({cur})"],
            "rows": [[v.get("name"), v.get("size"), v.get("state"),
                      v.get("cpu_avg"), v.get("cpu_peak"), v.get("cost")]
                     for v in (rs.get("vms") or [])],
            "money_col": 6,
        })

    esu = live.get("esu")
    if esu:
        out.append({
            "key": "esu",
            "title": "End of support",
            "facts": [f"{esu.get('out_of_support', 0)} out of support",
                      f"{esu.get('ending_soon', 0)} ending soon",
                      f"{esu.get('exposed', 0)} uncovered",
                      f"estimated {esu.get('estimated_monthly_cost') or 0:,.2f} {cur}/month"],
            "columns": ["Name", "Product", "Where", "Support ended", "ESU position",
                        f"Est. per month ({cur})"],
            "rows": [[m.get("name"), m.get("product"), m.get("kind"), m.get("support_ends"),
                      m.get("coverage"), m.get("monthly_esu_cost")]
                     for m in (esu.get("machines") or [])],
            "money_col": 6,
        })

    advisor = live.get("advisor")
    if advisor:
        out.append({
            "key": "advisor",
            "title": "Advisor recommendations",
            "facts": [f"{advisor.get('count', 0)} recommendation(s)",
                      f"estimated {advisor.get('estimated_annual_savings') or 0:,.2f} "
                      f"{cur}/year"],
            "columns": ["Impact", "Recommendation", "Resource", f"Annual saving ({cur})"],
            "rows": [[r.get("impact"), r.get("problem") or r.get("solution"),
                      r.get("resource"), r.get("annual_savings")]
                     for r in (advisor.get("recommendations") or [])],
            "money_col": 4,
        })

    return [t for t in out if t["rows"]]


# --------------------------------------------------------------------------- xlsx
def to_xlsx(data: dict[str, Any]) -> bytes:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, PieChart, Reference
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    cur = data["currency"]
    money = _money_format(cur)
    head_font = Font(bold=True, color="FFFFFF", size=10)
    head_fill = PatternFill("solid", fgColor="1F2933")
    title_font = Font(bold=True, size=14)

    def write_table(ws, rows: list[list], headers: list[str], start: int = 1,
                    money_cols: tuple[int, ...] = ()) -> int:
        for c, h in enumerate(headers, start=1):
            cell = ws.cell(row=start, column=c, value=h)
            cell.font = head_font
            cell.fill = head_fill
            cell.alignment = Alignment(horizontal="left")
        for r, row in enumerate(rows, start=start + 1):
            for c, v in enumerate(row, start=1):
                cell = ws.cell(row=r, column=c, value=v)
                if c in money_cols:
                    cell.number_format = money
        # Width from content: a column of "###" helps nobody.
        for c, h in enumerate(headers, start=1):
            longest = max([len(str(h))] + [len(str(row[c - 1])) for row in rows[:200]] or [8])
            ws.column_dimensions[get_column_letter(c)].width = min(max(longest + 3, 11), 46)
        ws.freeze_panes = ws.cell(row=start + 1, column=1)
        return start + len(rows) + 1

    # --- Summary
    ws = wb.active
    ws.title = "Summary"
    ws["A1"] = "CloudLens — Azure cost report"
    ws["A1"].font = title_font
    meta = [
        ("Generated", data["generated"]),
        ("Period", f"Last {data['days']} days"),
        ("Data through", data["as_of"] or "—"),
        ("Currency", cur),
        ("Scope", f"{data['subscription_count']} subscription(s)" if data["scoped"]
                  else "All accessible subscriptions"),
    ]
    for i, (k, v) in enumerate(meta, start=3):
        ws.cell(row=i, column=1, value=k).font = Font(bold=True)
        ws.cell(row=i, column=2, value=v)
    ws.cell(row=9, column=1, value="Total spend").font = Font(bold=True)
    total_cell = ws.cell(row=9, column=2, value=data["total"])
    total_cell.number_format = money
    total_cell.font = Font(bold=True, size=12)

    rows = [[a["label"], a["cost"],
             (a["cost"] / data["total"]) if data["total"] else 0] for a in data["areas"]]
    end = write_table(ws, rows, ["Area", f"Cost ({cur})", "Share"], start=11, money_cols=(2,))
    for r in range(12, 12 + len(rows)):
        ws.cell(row=r, column=3).number_format = "0.0%"

    if rows:
        pie = PieChart()
        pie.title = f"Spend by area — last {data['days']} days"
        pie.height, pie.width = 8, 13
        pie.add_data(Reference(ws, min_col=2, min_row=11, max_row=11 + len(rows)), titles_from_data=True)
        pie.set_categories(Reference(ws, min_col=1, min_row=12, max_row=11 + len(rows)))
        ws.add_chart(pie, f"E11")

    # --- One sheet per area
    for section in data["sections"]:
        name = section["label"][:28].replace("/", "-").replace("&", "and")
        sheet = wb.create_sheet(name)
        k = section["kpis"]
        sheet["A1"] = section["label"]
        sheet["A1"].font = title_font
        for i, (label, value, fmt) in enumerate([
            ("Spend", k["total"], money),
            ("Previous period", k["previous"], money),
            ("Change", (k["change_pct"] / 100) if k["change_pct"] is not None else None, "0.0%"),
            ("Share of total", (k["share_pct"] / 100) if k["share_pct"] is not None else None, "0.0%"),
            ("Billed resources", k["resources"], None),
            ("Services", k["services"], None),
        ], start=3):
            sheet.cell(row=i, column=1, value=label).font = Font(bold=True)
            cell = sheet.cell(row=i, column=2, value=value)
            if fmt:
                cell.number_format = fmt

        at = write_table(sheet, [[s["name"], s["cost"]] for s in section["services"]],
                         ["Service", f"Cost ({cur})"], start=10, money_cols=(2,))
        at = write_table(sheet, [[r["name"], r["cost"]] for r in section["regions"]],
                         ["Region", f"Cost ({cur})"], start=at + 1, money_cols=(2,))
        write_table(sheet,
                    [[r["name"], r["grp"], r["service"], r["region"], r["cost"]]
                     for r in section["resources_top"]],
                    ["Resource", "Resource group", "Service", "Region", f"Cost ({cur})"],
                    start=at + 1, money_cols=(5,))

        trend = section.get("trend") or {}
        if trend.get("labels"):
            tsheet = wb.create_sheet(f"{name[:24]} trend")
            write_table(tsheet, list(zip(trend["labels"], trend["values"])),
                        ["Date", f"Cost ({cur})"], money_cols=(2,))
            chart = BarChart()
            chart.title = f"{section['label']} — daily spend"
            chart.height, chart.width = 8, 20
            chart.add_data(Reference(tsheet, min_col=2, min_row=1, max_row=1 + len(trend["labels"])),
                           titles_from_data=True)
            chart.set_categories(Reference(tsheet, min_col=1, min_row=2,
                                           max_row=1 + len(trend["labels"])))
            tsheet.add_chart(chart, "D2")

    # --- Live findings, one sheet each
    for t in live_tables(data):
        sheet = wb.create_sheet(t["title"][:28].replace("/", "-").replace("&", "and"))
        sheet["A1"] = t["title"]
        sheet["A1"].font = title_font
        sheet["A2"] = " · ".join(t["facts"])
        sheet["A2"].font = Font(color="626A76", size=10)
        write_table(sheet, t["rows"], t["columns"], start=4, money_cols=(t["money_col"],))

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# --------------------------------------------------------------------------- pptx
def to_pptx(data: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)  # 16:9
    blank = prs.slide_layouts[6]
    cur = data["currency"]
    symbol = fx_symbol(cur)
    ink = RGBColor(0x20, 0x1F, 0x1E)      # Fluent neutral, as ink
    dim = RGBColor(0x60, 0x5E, 0x5C)      # Fluent secondary text
    accent = RGBColor(0x00, 0x78, 0xD4)   # the reference blue

    def textbox(slide, text, left, top, width, height, size=18, bold=False, colour=ink):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = "Segoe UI"
        return box

    def money(v):
        return f"{symbol}{v:,.2f}" if v is not None else "—"

    # --- Title
    s = prs.slides.add_slide(blank)
    textbox(s, "CloudLens", Inches(.9), Inches(2.0), Inches(11), Inches(1), 44, True)
    textbox(s, "SEE MORE.  OPTIMIZE BETTER.", Inches(.95), Inches(2.95), Inches(11),
            Inches(.45), 11, False, dim)
    textbox(s, "Azure cost report", Inches(.9), Inches(3.55), Inches(11), Inches(.6), 22,
            False, ink)
    scope = (f"{data['subscription_count']} subscription(s)" if data["scoped"]
             else "All accessible subscriptions")
    textbox(s, f"Last {data['days']} days · data through {data['as_of'] or '—'} · {scope}",
            Inches(.9), Inches(4.2), Inches(11), Inches(.5), 15, False, dim)
    textbox(s, f"Total {money(data['total'])} {cur}",
            Inches(.9), Inches(4.85), Inches(11), Inches(.8), 26, True, accent)
    textbox(s, f"Generated {data['generated']}",
            Inches(.9), Inches(6.5), Inches(11), Inches(.4), 11, False, dim)

    # --- Where the money goes: chart beside the numbers, not one after the other.
    if data["areas"]:
        s = prs.slides.add_slide(blank)
        textbox(s, "Where the money goes", Inches(.6), Inches(.4), Inches(9), Inches(.7), 28, True)

        chart_data = CategoryChartData()
        chart_data.categories = [a["label"] for a in data["areas"]]
        chart_data.add_series("Cost", [a["cost"] for a in data["areas"]])
        # A native chart, not an image: the deck stays editable in PowerPoint.
        frame = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(.5), Inches(1.3),
                                   Inches(7), Inches(5.6), chart_data)
        chart = frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        chart.font.size = Pt(12)

        rows, cols = len(data["areas"]) + 1, 3
        table = s.shapes.add_table(rows, cols, Inches(7.9), Inches(1.4),
                                   Inches(4.9), Inches(.4 * rows)).table
        for i, h in enumerate(["Area", "Cost", "Share"]):
            cell = table.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
        for r, a in enumerate(data["areas"], start=1):
            share = (a["cost"] / data["total"] * 100) if data["total"] else 0
            for c, v in enumerate([a["label"], money(a["cost"]), f"{share:.1f}%"]):
                cell = table.cell(r, c)
                cell.text = str(v)
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)

    # --- One slide per area, biggest first
    for section in sorted(data["sections"], key=lambda x: -(x["kpis"]["total"] or 0))[:8]:
        k = section["kpis"]
        s = prs.slides.add_slide(blank)
        textbox(s, section["label"], Inches(.6), Inches(.4), Inches(9), Inches(.7), 28, True)

        change = (f"{k['change_pct']:+.1f}% vs previous period"
                  if k["change_pct"] is not None else "no prior period to compare")
        textbox(s, f"{money(k['total'])}   ·   {change}   ·   {k['resources']} resources",
                Inches(.6), Inches(1.15), Inches(12), Inches(.5), 15, False, dim)

        trend = section.get("trend") or {}
        if len(trend.get("labels", [])) > 1:
            cd = CategoryChartData()
            cd.categories = trend["labels"]
            cd.add_series(section["label"], trend["values"])
            frame = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(.5), Inches(1.8),
                                       Inches(7.4), Inches(5), cd)
            frame.chart.has_legend = False
            frame.chart.font.size = Pt(10)

        top = section["services"][:10]
        if top:
            rows = len(top) + 1
            table = s.shapes.add_table(rows, 2, Inches(8.2), Inches(1.9),
                                       Inches(4.6), Inches(.35 * rows)).table
            for i, h in enumerate(["Top services", "Cost"]):
                cell = table.cell(0, i)
                cell.text = h
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
            for r, item in enumerate(top, start=1):
                for c, v in enumerate([item["name"], money(item["cost"])]):
                    cell = table.cell(r, c)
                    cell.text = str(v)
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10)

    # --- Live findings, one slide each
    for t in live_tables(data):
        s = prs.slides.add_slide(blank)
        textbox(s, t["title"], Inches(.6), Inches(.4), Inches(11), Inches(.7), 28, True)
        textbox(s, "   ·   ".join(t["facts"]), Inches(.6), Inches(1.15), Inches(12),
                Inches(.5), 15, False, dim)

        shown = t["rows"][:12]
        rows, cols = len(shown) + 1, len(t["columns"])
        table = s.shapes.add_table(rows, cols, Inches(.6), Inches(1.9),
                                   Inches(12.1), Inches(min(.34 * rows, 5))).table
        for i, h in enumerate(t["columns"]):
            cell = table.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
        for r, row in enumerate(shown, start=1):
            for c, v in enumerate(row):
                cell = table.cell(r, c)
                cell.text = money(v) if (c + 1 == t["money_col"] and isinstance(v, (int, float))) \
                    else ("—" if v is None else str(v))
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(9.5)
        if len(t["rows"]) > len(shown):
            textbox(s, f"Showing {len(shown)} of {len(t['rows'])} — the full list is in the "
                       "Excel and CSV exports.",
                    Inches(.6), Inches(7), Inches(12), Inches(.4), 10, False, dim)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# ------------------------------------------------------------------ csv / md / json
def to_csv(data: dict[str, Any]) -> bytes:
    """One flat table. A report split across sheets is useless to a script."""
    out = io.StringIO()
    w = csv.writer(out, lineterminator="\n")
    cur = data["currency"]
    w.writerow(["Area", "Type", "Name", "Resource group", "Service", "Region", f"Cost ({cur})"])
    for a in data["areas"]:
        w.writerow([a["label"], "area total", "", "", "", "", f"{a['cost']:.2f}"])
    for section in data["sections"]:
        for s in section["services"]:
            w.writerow([section["label"], "service", s["name"], "", "", "", f"{s['cost']:.2f}"])
        for r in section["regions"]:
            w.writerow([section["label"], "region", r["name"], "", "", "", f"{r['cost']:.2f}"])
        for r in section["resources_top"]:
            w.writerow([section["label"], "resource", r["name"], r["grp"], r["service"],
                        r["region"], f"{r['cost']:.2f}"])
    for t in live_tables(data):
        for row in t["rows"]:
            cost = row[t["money_col"] - 1]
            w.writerow([t["title"], t["key"], str(row[0] or ""), "", "", "",
                        f"{cost:.2f}" if isinstance(cost, (int, float)) else ""])
    return out.getvalue().encode("utf-8-sig")  # BOM so Excel opens UTF-8 correctly


def to_json(data: dict[str, Any]) -> bytes:
    """The whole selection, unflattened — for a script, a notebook or a pipeline."""
    import json

    return json.dumps(data, indent=2, default=str).encode("utf-8")


def to_md(data: dict[str, Any]) -> bytes:
    cur = data["currency"]
    symbol = fx_symbol(cur)
    money = lambda v: f"{symbol}{v:,.2f}"

    lines = [
        "# CloudLens — Azure cost report",
        "",
        f"- **Period**: last {data['days']} days (data through {data['as_of'] or '—'})",
        f"- **Scope**: {f'{data['subscription_count']} subscription(s)' if data['scoped'] else 'all accessible subscriptions'}",
        f"- **Total**: {money(data['total'])} {cur}",
        f"- **Generated**: {data['generated']}",
        "",
        "## Where the money goes",
        "",
        "| Area | Cost | Share |",
        "|---|---:|---:|",
    ]
    for a in data["areas"]:
        share = (a["cost"] / data["total"] * 100) if data["total"] else 0
        lines.append(f"| {a['label']} | {money(a['cost'])} | {share:.1f}% |")

    for section in data["sections"]:
        k = section["kpis"]
        change = f"{k['change_pct']:+.1f}%" if k["change_pct"] is not None else "—"
        lines += [
            "", f"## {section['label']}", "",
            f"{money(k['total'])} · {change} vs previous period · "
            f"{k['resources']} resources · {k['services']} services", "",
            "| Service | Cost |", "|---|---:|",
        ]
        lines += [f"| {s['name']} | {money(s['cost'])} |" for s in section["services"][:10]]
        if section["resources_top"]:
            lines += ["", "| Resource | Resource group | Service | Cost |", "|---|---|---|---:|"]
            lines += [f"| {r['name']} | {r['grp']} | {r['service']} | {money(r['cost'])} |"
                      for r in section["resources_top"][:10]]

    for t in live_tables(data):
        lines += ["", f"## {t['title']}", "", " · ".join(t["facts"]), "",
                  "| " + " | ".join(t["columns"]) + " |",
                  "|" + "---|" * len(t["columns"])]
        for row in t["rows"][:20]:
            cells = [money(v) if (i + 1 == t["money_col"] and isinstance(v, (int, float)))
                     else ("—" if v is None else str(v)) for i, v in enumerate(row)]
            lines.append("| " + " | ".join(cells) + " |")

    return ("\n".join(lines) + "\n").encode("utf-8")


BUILDERS = {"xlsx": to_xlsx, "pptx": to_pptx, "csv": to_csv, "md": to_md, "json": to_json}


def build(fmt: str, scope: list[str] | None, days: int = 30,
          selection: dict[str, Any] | None = None,
          live_data: dict[str, Any] | None = None,
          currency: str | None = None) -> tuple[bytes, str, str]:
    """Returns (bytes, filename, media type)."""
    if fmt not in BUILDERS:
        raise KeyError(fmt)
    data = collect(scope, days=days, selection=selection, live_data=live_data,
                   currency=currency)
    payload = BUILDERS[fmt](data)
    stamp = datetime.now(timezone.utc).strftime("%Y%m%d")
    # The currency goes in the filename when it is not the billed one. Two files a week apart,
    # one in dollars and one in rupees, are otherwise indistinguishable in a downloads folder.
    tag = f"-{currency.lower()}" if currency and currency.upper() != "BILLED" else ""
    return payload, f"cloudlens-cost-report-{stamp}{tag}.{fmt}", FORMATS[fmt][1]
